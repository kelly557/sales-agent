#!/usr/bin/env python3
import argparse
import json
import os
import sys
from typing import Any

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    print("python-pptx is required. Install with: python3 -m pip install python-pptx", file=sys.stderr)
    sys.exit(2)


def text(value: Any, default: str = "") -> str:
    return value.strip() if isinstance(value, str) and value.strip() else default


def rows(value: Any) -> list[dict[str, Any]]:
    return value if isinstance(value, list) else []


def add_title(slide, title: str, subtitle: str = ""):
    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.45), Inches(12.1), Inches(0.72))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 41, 59)

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.67), Inches(1.12), Inches(12), Inches(0.35))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(100, 116, 139)


def add_footer(slide, page: int):
    footer = slide.shapes.add_textbox(Inches(0.65), Inches(7.05), Inches(12), Inches(0.25))
    frame = footer.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = f"DocTech 定制方案生成 | {page:02d}"
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(148, 163, 184)


def add_bullets(slide, items: list[str], x: float, y: float, w: float, h: float, font_size: int = 15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()

    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(51, 65, 85)
        p.space_after = Pt(8)


def add_table(slide, columns: list[str], data: list[dict[str, Any]], x: float, y: float, w: float, h: float):
    if not data:
        add_bullets(slide, ["暂无数据"], x, y, w, h)
        return

    table_shape = slide.shapes.add_table(len(data) + 1, len(columns), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table

    for col_index, column in enumerate(columns):
        cell = table.cell(0, col_index)
        cell.text = column
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(15, 118, 110)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(10)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)

    for row_index, row in enumerate(data, 1):
        for col_index, column in enumerate(columns):
            cell = table.cell(row_index, col_index)
            cell.text = text(row.get(column), "-")
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = RGBColor(51, 65, 85)


def normalize_slides(payload: dict[str, Any]) -> list[dict[str, str]]:
    model_slides = rows(payload.get("modelSlides"))
    fallback_slides = rows(payload.get("slides"))
    source = model_slides if model_slides else fallback_slides

    normalized = []
    for item in source:
        if not isinstance(item, dict):
            continue
        title = text(item.get("title"))
        content = text(item.get("content"))
        if title and content:
            normalized.append({"title": title, "content": content})

    if normalized:
        return normalized

    return [
        {"title": "项目背景与客户诉求", "content": "说明项目范围、客户关注点和方案汇报目标。"},
        {"title": "推荐方案", "content": "展示体系选择、施工工艺、质量控制和风险提示。"},
        {"title": "下一步计划", "content": "汇总待复核资料、审批节点和交付清单。"},
    ]


def build_ppt(payload: dict[str, Any], output_path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    project = text(payload.get("projectName"), "定制方案 PPT")
    customer = text(payload.get("customer"), "客户")
    scenario = text(payload.get("scenario"), "项目场景")
    demand = text(payload.get("demand"), "客户诉求")
    slides = normalize_slides(payload)

    cover = prs.slides.add_slide(blank)
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
    add_title(cover, project, f"{customer} | {scenario}")
    add_bullets(
        cover,
        [
            f"项目位置：{text(payload.get('location'), '-')}",
            f"涂装面积：{text(payload.get('area'), '-')}",
            f"推荐体系：{text(payload.get('coatingSystem'), '-')}",
            f"客户诉求：{demand}",
        ],
        0.9,
        2.0,
        11.6,
        2.8,
        18,
    )
    add_footer(cover, 1)

    agenda = prs.slides.add_slide(blank)
    add_title(agenda, "目录", "根据项目参数、知识库、规则校验和审批要求生成")
    add_bullets(agenda, [f"{index + 1:02d}. {item['title']}" for index, item in enumerate(slides)], 0.9, 1.6, 11.5, 5.1, 14)
    add_footer(agenda, 2)

    page = 3
    for index, item in enumerate(slides, 1):
        slide = prs.slides.add_slide(blank)
        add_title(slide, item["title"], f"方案内容页 {index:02d}")
        add_bullets(slide, split_content(item["content"]), 0.9, 1.65, 11.4, 4.95, 17)
        add_footer(slide, page)
        page += 1

    knowledge = prs.slides.add_slide(blank)
    add_title(knowledge, "知识引用", "生成内容使用的企业知识、产品资料与规则来源")
    add_table(knowledge, ["name", "type", "status", "owner"], rows(payload.get("knowledge")), 0.6, 1.55, 12.0, 4.9)
    add_footer(knowledge, page)
    page += 1

    rules = prs.slides.add_slide(blank)
    add_title(rules, "规则校验", "体系选型、用量估算、施工边界和风险提示")
    add_table(rules, ["rule", "result", "detail"], rows(payload.get("rules")), 0.6, 1.55, 12.0, 4.9)
    add_footer(rules, page)
    page += 1

    citations = prs.slides.add_slide(blank)
    add_title(citations, "内容溯源", "关键页内容对应的资料来源与证据")
    add_table(citations, ["slide", "source", "evidence"], rows(payload.get("citations")), 0.6, 1.55, 12.0, 4.9)
    add_footer(citations, page)
    page += 1

    approvals = prs.slides.add_slide(blank)
    add_title(approvals, "审批与交付", "复核节点、负责人和交付关注点")
    add_table(approvals, ["node", "assignee", "status", "focus"], rows(payload.get("approvals")), 0.6, 1.55, 12.0, 4.9)
    add_footer(approvals, page)

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)


def split_content(content: str) -> list[str]:
    parts = [part.strip(" -•\t") for part in content.replace("；", "\n").replace("。", "。\n").splitlines()]
    return [part for part in parts if part][:6] or [content]


def main():
    parser = argparse.ArgumentParser(description="Generate editable proposal PPTX locally.")
    parser.add_argument("--input", required=True, help="Path to proposal JSON payload")
    parser.add_argument("--output", required=True, help="Output .pptx path")
    args = parser.parse_args()

    with open(args.input, "r", encoding="utf-8") as f:
        payload = json.load(f)

    build_ppt(payload, args.output)
    print(f"DONE saved={os.path.abspath(args.output)}")


if __name__ == "__main__":
    main()
